import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
import pytesseract
import io


# Enter the path to get the input feeded to the api
PATH = "./NoogatAssignment.pptx"


def extract_text_from_shapes(slide):
    """Extract visible text from slide shapes."""
    text_runs = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            text = " ".join(run.text.strip() for run in paragraph.runs if run.text.strip())
            if text:
                text_runs.append(text)
    return text_runs

def extract_text_from_images(slide):
    """Extract text from images in a slide using OCR."""
    image_texts = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            image = shape.image
            # Save image temporarily in memory
            image_bytes = io.BytesIO(image.blob)
            img = Image.open(image_bytes)
            text = pytesseract.image_to_string(img)
            if text.strip():
                image_texts.append(text.strip())
    return image_texts

def extract_chart_data(slide):
    """Extract data from charts in a slide, safely handling missing attributes."""
    chart_data_list = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.CHART:
            chart = shape.chart
            chart_info = {"chart_type": str(chart.chart_type), "series": []}
            for series in chart.series:
                # Some chart types may not have categories
                categories = []
                if hasattr(series, "categories") and series.categories is not None:
                    categories = [str(c) for c in series.categories]
                
                values = []
                if hasattr(series, "values") and series.values is not None:
                    values = [v for v in series.values]
                
                chart_info["series"].append({
                    "name": series.name if hasattr(series, "name") else "Unnamed Series",
                    "values": values,
                    "categories": categories
                })
            chart_data_list.append(chart_info)
    return chart_data_list



def process_pptx(file_path):
    prs = Presentation(file_path)
    all_data = []

    for i, slide in enumerate(prs.slides, start=1):
        slide_data = {
            "slide_number": i,
            "text": extract_text_from_shapes(slide),
            "image_text": extract_text_from_images(slide),
            "charts": extract_chart_data(slide)
        }
        all_data.append(slide_data)

    return all_data

if __name__ == "__main__":
    data = process_pptx(PATH)


    for slide in data:
        print(f"\n--- Slide {slide['slide_number']} ---")
        print("Text from shapes:")
        for t in slide['text']:
            print(f"  - {t}")
        print("Text from images (OCR):")
        for t in slide['image_text']:
            print(f"  - {t}")
        print("Chart data:")
        for chart in slide['charts']:
            print(f"  Chart type: {chart['chart_type']}")
            for series in chart['series']:
                print(f"    Series: {series['name']}")
                print(f"    Categories: {series['categories']}")
                print(f"    Values: {series['values']}")

